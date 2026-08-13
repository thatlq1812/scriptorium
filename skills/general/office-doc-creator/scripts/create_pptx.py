#!/usr/bin/env python3
"""Create a .pptx from a simple JSON content spec.

Usage:
    python create_pptx.py <content.json> <output.pptx>

content.json shape:
{
  "slides": [
    {"title": "Slide 1", "bullets": ["point one", "point two"]}
  ]
}
"""
from __future__ import annotations

import argparse
import json
from pathlib import Path

from pptx import Presentation


def build(content: dict, output_path: Path) -> None:
    prs = Presentation()
    title_and_content = prs.slide_layouts[1]

    for slide_spec in content.get("slides", []):
        slide = prs.slides.add_slide(title_and_content)
        slide.shapes.title.text = slide_spec.get("title", "")
        body = slide.placeholders[1].text_frame
        bullets = slide_spec.get("bullets", [])
        if bullets:
            body.text = bullets[0]
            for bullet in bullets[1:]:
                p = body.add_paragraph()
                p.text = bullet

    prs.save(str(output_path))


def main() -> int:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("content_json", type=Path)
    parser.add_argument("output_pptx", type=Path)
    args = parser.parse_args()

    content = json.loads(args.content_json.read_text(encoding="utf-8"))
    build(content, args.output_pptx)
    print(f"OK: wrote {args.output_pptx}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
