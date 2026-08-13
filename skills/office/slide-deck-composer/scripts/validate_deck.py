#!/usr/bin/env python3
"""
Real structural validator for a compiled slide-deck-composer output
file. Every check here is grounded in an actual failure mode this
skill's own build/fix history documented (see SKILL.md "Verified"
bugs #1-10) -- not a generic OOXML linter guessing at what might go
wrong.

Usage:
    python validate_deck.py <output.pptx> [--original template.pptx]

Checks (always run):
  1. Valid zip/OPC package -- readable zip, zero corrupt members
     (``zipfile.testzip()``).
  2. Every slide part's XML is well-formed (``lxml.etree.parse``) --
     guards against the class of bug where a script mutates a
     ``_txBody``/``_element`` tree into something that LOOKS right in
     python-pptx's in-memory object graph but serializes to broken
     XML (python-pptx itself won't catch this at save-time).
  3. Every slide part's XML is VALID against the real ECMA-376
     PresentationML schema (v0.7.0, ``ooxml_schema_valid``) -- a
     stronger, categorically different check than #2: well-formed XML
     can still violate the schema (wrong element order, a required
     attribute missing, an out-of-enum value) in ways PowerPoint may
     silently "fix"/reinterpret or outright refuse to open, and #2
     alone cannot catch any of that. Uses ``assets/xsd/pml.xsd`` (the
     real ECMA-376 schema set, copied verbatim from ``python-pptx``'s
     own repo under the same MIT license this skill's core dependency
     already carries -- see ``assets/xsd/NOTICE.md``) via
     ``lxml.etree.XMLSchema``, no new dependency. Feasibility (does a
     real-world OOXML producer's output actually pass strict schema
     validation, or is this too strict to be useful against real
     Slidesgo/Google-Slides-exported templates) was confirmed
     empirically before adoption, not assumed -- see SKILL.md
     "Verified" v0.7.0.
  4. No dangling relationship references -- every ``r:id``/``r:embed``/
     ``r:link`` attribute actually used inside a slide's XML resolves
     to a real relationship in that slide's own ``.rels`` part. A
     stale/removed relationship reference here is exactly the class of
     bug ``image_injector.py``'s blip-swap or ``slide_duplicator.py``'s
     rId-remap could silently produce if either had an off-by-one.
  5. No leaked vendor instructional/watermark text -- reuses this
     skill's own real, tested detection strings
     (``slide_classifier._INSTRUCTION_MARKERS`` /
     ``shape_roles.WATERMARK_TEXT_PATTERNS``) as a REGRESSION check
     against bugs #4 and #10 (a Portuguese "About this template" page
     and a "Fonts & colors used" credits page both leaked into real
     output before those fixes) -- if either pattern set ever regresses,
     this check catches it mechanically instead of requiring another
     manual visual-review round to notice.
  6. Slide dimensions match the compiled output's own declared
     presentation size (sanity check that nothing corrupted
     ``p:sldSz`` during compile).

Checks (only with --original, informational/baseline-only, never a
hard fail -- these exist so a TEMPLATE's own pre-existing quirks
aren't mistaken for a defect this skill introduced):
  7. Output slide dimensions vs. the original template's own
     dimensions (hard fail if they differ -- this skill's compile
     pipeline never intentionally resizes a slide).
  8. Theme major/minor font: reports whether the output's theme fonts
     differ from the original template's (informational only --
     ``--heading-font``/``--body-font`` deliberately change this, so a
     difference is not itself a defect).

Exit codes: 0 = all hard checks passed, 2 = malformed/missing input,
3 = one or more hard checks failed (see JSON report's "failures").
"""
from __future__ import annotations

import argparse
import json
import sys
import zipfile
from pathlib import Path

from lxml import etree
from pptx import Presentation

sys.path.insert(0, str(Path(__file__).resolve().parent))
import slide_classifier  # noqa: E402
import shape_roles  # noqa: E402

for _stream in (sys.stdout, sys.stderr):
    if hasattr(_stream, "reconfigure"):
        _stream.reconfigure(encoding="utf-8")

_NS_R = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
_RID_ATTRS = (f"{{{_NS_R}}}id", f"{{{_NS_R}}}embed", f"{{{_NS_R}}}link")

_XSD_DIR = Path(__file__).resolve().parent.parent / "assets" / "xsd"
_PML_XSD_PATH = _XSD_DIR / "pml.xsd"
_pml_schema_cache: "etree.XMLSchema | None" = None


def _load_pml_schema() -> etree.XMLSchema:
    """Lazily load and cache the real ECMA-376 PresentationML schema
    (see assets/xsd/NOTICE.md for provenance/license). Loading parses
    24 interconnected .xsd files (pml.xsd xsd:include/import-s the
    DrawingML and shared-type schemas also in that directory) -- cached
    at module level so validating many slides in one run only pays
    this cost once."""
    global _pml_schema_cache
    if _pml_schema_cache is None:
        if not _PML_XSD_PATH.is_file():
            raise ValidateDeckError(
                f"OOXML schema not found at {_PML_XSD_PATH} -- assets/xsd/ appears to be "
                f"missing or incomplete, cannot run the schema-validity check"
            )
        _pml_schema_cache = etree.XMLSchema(etree.parse(str(_PML_XSD_PATH)))
    return _pml_schema_cache


class ValidateDeckError(Exception):
    """Raised for any refuse-loudly condition (malformed/missing input)."""


def check_zip_integrity(path: Path) -> dict:
    if not zipfile.is_zipfile(path):
        return {"check": "zip_integrity", "ok": False, "detail": "not a valid zip/OPC package"}
    with zipfile.ZipFile(path) as zf:
        bad = zf.testzip()
    return {
        "check": "zip_integrity", "ok": bad is None,
        "detail": f"corrupt member: {bad}" if bad else "clean",
    }


def _slide_part_names(zf: zipfile.ZipFile) -> list[str]:
    return sorted(
        n for n in zf.namelist()
        if n.startswith("ppt/slides/slide") and n.endswith(".xml")
    )


def check_slide_xml_wellformed(path: Path) -> dict:
    failures: list[str] = []
    with zipfile.ZipFile(path) as zf:
        for name in _slide_part_names(zf):
            try:
                etree.fromstring(zf.read(name))
            except etree.XMLSyntaxError as e:
                failures.append(f"{name}: {e}")
    return {
        "check": "slide_xml_wellformed", "ok": not failures,
        "detail": failures or "all slide parts well-formed XML",
    }


def check_ooxml_schema_valid(path: Path) -> dict:
    """Validate every slide part against the real ECMA-376
    PresentationML schema (v0.7.0) -- see module docstring check #3
    and assets/xsd/NOTICE.md. A categorically stronger check than
    well-formedness alone: well-formed XML can still violate the
    schema (wrong element order, a required attribute missing, an
    out-of-enum value)."""
    schema = _load_pml_schema()
    failures: list[str] = []
    with zipfile.ZipFile(path) as zf:
        for name in _slide_part_names(zf):
            try:
                doc = etree.fromstring(zf.read(name))
            except etree.XMLSyntaxError:
                # Already reported by check_slide_xml_wellformed --
                # don't double-report the same defect under a
                # different check name, and schema-validating
                # unparseable XML isn't meaningful anyway.
                continue
            if not schema.validate(doc):
                for error in schema.error_log:
                    failures.append(f"{name}:{error.line}: {error.message}")
    return {
        "check": "ooxml_schema_valid", "ok": not failures,
        "detail": failures or f"all slide parts valid against ECMA-376 PresentationML schema ({_PML_XSD_PATH.name})",
    }


def check_no_dangling_relationships(path: Path) -> dict:
    failures: list[str] = []
    with zipfile.ZipFile(path) as zf:
        names = set(zf.namelist())
        for slide_name in _slide_part_names(zf):
            slide_dir, slide_file = slide_name.rsplit("/", 1)
            rels_name = f"{slide_dir}/_rels/{slide_file}.rels"
            declared_rids: set[str] = set()
            if rels_name in names:
                rels_root = etree.fromstring(zf.read(rels_name))
                for rel in rels_root:
                    rid = rel.get("Id")
                    if rid:
                        declared_rids.add(rid)

            slide_root = etree.fromstring(zf.read(slide_name))
            used_rids: set[str] = set()
            for el in slide_root.iter():
                for attr in _RID_ATTRS:
                    val = el.get(attr)
                    if val:
                        used_rids.add(val)

            missing = used_rids - declared_rids
            if missing:
                failures.append(f"{slide_name}: r:id(s) referenced but not declared in {rels_name}: {sorted(missing)}")
    return {
        "check": "no_dangling_relationships", "ok": not failures,
        "detail": failures or "every used r:id/r:embed/r:link resolves to a declared relationship",
    }


def check_no_leaked_vendor_text(path: Path) -> dict:
    """Regression check for SKILL.md bugs #4 and #10: reuses this
    skill's own real detection heuristics against the COMPILED OUTPUT
    (not the template) -- a genuine defect would mean instructional/
    credits content from the template slipped through uncleaned into
    real generated slides."""
    prs = Presentation(str(path))
    failures: list[str] = []
    for i, slide in enumerate(prs.slides):
        try:
            if slide_classifier.is_instruction_slide(slide):
                failures.append(f"output slide {i}: matches is_instruction_slide (vendor/credits page leaked into real output)")
        except Exception as e:
            failures.append(f"output slide {i}: is_instruction_slide check raised {e!r}")
        for shape in slide.shapes:
            try:
                if not shape.has_text_frame:
                    continue
                text = shape.text_frame.text or ""
            except Exception:
                continue
            if len(text) > slide_classifier._INSTRUCTION_PARAGRAPH_MIN_CHARS and shape_roles.is_watermark_text(text):
                failures.append(
                    f"output slide {i}, shape {getattr(shape, 'shape_id', '?')}: "
                    f"long text ({len(text)} chars) matches watermark/vendor-brand "
                    f"pattern -- likely leaked instructional content, not a "
                    f"legitimate short attribution footer"
                )
            if shape_roles.is_vendor_instructional_caption(text):
                failures.append(
                    f"output slide {i}, shape {getattr(shape, 'shape_id', '?')}: "
                    f"text matches a vendor chart-editing caption pattern -- "
                    f"leaked instructional filler, not real content"
                )
    return {
        "check": "no_leaked_vendor_text", "ok": not failures,
        "detail": failures or "no vendor instructional/credits text detected in output",
    }


def check_slide_dimensions_self_consistent(path: Path) -> dict:
    prs = Presentation(str(path))
    declared_w, declared_h = prs.slide_width, prs.slide_height
    failures: list[str] = []
    if not declared_w or not declared_h:
        failures.append("presentation.xml has no usable p:sldSz (slide_width/slide_height missing)")
    return {
        "check": "slide_dimensions_self_consistent", "ok": not failures,
        "detail": failures or f"slide size {declared_w}x{declared_h} EMU, {len(prs.slides)} slide(s)",
    }


def check_against_original(output_path: Path, original_path: Path) -> list[dict]:
    out_prs = Presentation(str(output_path))
    orig_prs = Presentation(str(original_path))
    results: list[dict] = []

    dims_match = (out_prs.slide_width, out_prs.slide_height) == (orig_prs.slide_width, orig_prs.slide_height)
    results.append({
        "check": "slide_dimensions_match_original", "ok": dims_match,
        "detail": (
            f"output={out_prs.slide_width}x{out_prs.slide_height}, "
            f"original={orig_prs.slide_width}x{orig_prs.slide_height}"
        ),
    })

    import font_manager
    out_theme = font_manager.FontManager(out_prs).get_theme_fonts()
    orig_theme = font_manager.FontManager(orig_prs).get_theme_fonts()
    theme_matches = (out_theme.major_font == orig_theme.major_font
                      and out_theme.minor_font == orig_theme.minor_font)
    results.append({
        "check": "theme_fonts_vs_original", "ok": True,  # informational only, never a hard fail
        "informational": True,
        "detail": (
            f"output major/minor={out_theme.major_font}/{out_theme.minor_font}, "
            f"original major/minor={orig_theme.major_font}/{orig_theme.minor_font}, "
            f"unchanged={theme_matches} (a difference is EXPECTED and correct when "
            f"--heading-font/--body-font/--vietnamese was used -- not itself a defect)"
        ),
    })
    return results


def validate_deck(output_path: Path, original_path: Path | None) -> dict:
    if not output_path.is_file():
        raise ValidateDeckError(f"output file not found: {output_path}")
    if original_path is not None and not original_path.is_file():
        raise ValidateDeckError(f"--original template file not found: {original_path}")

    zip_check = check_zip_integrity(output_path)
    checks = [zip_check]
    if not zip_check["ok"]:
        # Every remaining check reads the zip/OPC package directly --
        # not worth running (and would just raise BadZipFile) once the
        # package itself is confirmed unreadable. Report this as a
        # real hard failure, not a script crash.
        return {
            "output": str(output_path),
            "original": str(original_path) if original_path else None,
            "all_checks_passed": False,
            "checks": checks,
        }

    checks.extend([
        check_slide_xml_wellformed(output_path),
        check_ooxml_schema_valid(output_path),
        check_no_dangling_relationships(output_path),
        check_no_leaked_vendor_text(output_path),
        check_slide_dimensions_self_consistent(output_path),
    ])
    if original_path is not None:
        checks.extend(check_against_original(output_path, original_path))

    hard_failures = [c for c in checks if not c["ok"] and not c.get("informational", False)]
    return {
        "output": str(output_path),
        "original": str(original_path) if original_path else None,
        "all_checks_passed": not hard_failures,
        "checks": checks,
    }


def main() -> int:
    parser = argparse.ArgumentParser(description="Validate a compiled slide-deck-composer output .pptx.")
    parser.add_argument("output", type=Path, help="Compiled output .pptx to validate")
    parser.add_argument(
        "--original", type=Path, default=None,
        help="Caller-supplied template .pptx used to compile this output (baseline sanity checks)",
    )
    args = parser.parse_args()

    try:
        result = validate_deck(args.output, args.original)
    except ValidateDeckError as e:
        print(f"ERROR: {e}", file=sys.stderr)
        return 2
    except Exception as e:  # noqa: BLE001 - top-level CLI boundary
        print(f"ERROR: unexpected failure: {e}", file=sys.stderr)
        return 2

    print(json.dumps(result, indent=2, ensure_ascii=False))
    return 0 if result["all_checks_passed"] else 3


if __name__ == "__main__":
    raise SystemExit(main())
