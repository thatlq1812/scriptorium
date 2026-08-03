#!/usr/bin/env python3
"""
slide-deck-composer entry point: clone a caller-supplied real .pptx
template's slides and inject caller-supplied content into them,
preserving the template's own design (fonts, colors, layout geometry,
decorative elements) instead of drawing generic geometric slides from
scratch.

Architecture: see SKILL.md "Architecture decision" section. This is
the "clone-and-inject" pipeline -- NOT a from-scratch layout builder.

Usage:
    python compile_deck.py <template.pptx> <content.json> <output.pptx>
        [--heading-font "Montserrat"] [--body-font "Open Sans"]
        [--vietnamese]

content.json shape:
    {
      "slides": [
        {"layout": "title", "markdown": "# Deck Title\\nA subtitle line"},
        {"layout": "standard", "markdown": "## Slide heading\\n- bullet one\\n- bullet two",
         "source_slide_idx": 4},
        {"layout": "two_column", "markdown": "## Compare\\n### Left\\n- a\\n### Right\\n- b"},
        {"layout": "image", "markdown": "## Caption title\\nCaption text",
         "image_path": "assets/photo.jpg"},
        {"layout": "quote", "markdown": "# A standalone statement"},
        {"layout": "section", "markdown": "## Section N\\nSection subtitle"},
        {"layout": "custom", "source_slide_idx": 10,
         "slots": [{"role": "title", "text": "$4.5M"},
                   {"role": "body", "text": "Budget for the new products/services"}]}
      ]
    }

``layout`` must be one of the 7 keys in constants.LAYOUT_NAME_MAP --
see SKILL.md "Known limitations" for the ~14 layout SHAPES not built
as a NAMED layout (use "custom" for those, see below); an unrecognized
name is refused with a clear error, never silently downgraded.

``layout: "custom"`` (no ``markdown`` field; ``source_slide_idx`` and
``slots`` instead): an ad-hoc slot mapping for a real template slide
whose shape doesn't match one of the 6 named layouts (a stat callout,
an org chart, a 3-column agenda, an icon grid, ...) -- see SKILL.md
"Custom layout". ``source_slide_idx`` is MANDATORY for this layout
(there's no generic pool to search automatically -- render
``thumbnail_grid.py``, look at the slide, count its real distinct text
areas, then declare that many slots). ``slots`` is a list of
``{"role": "title"|"body", "text": "..."|["...", "..."]}``: at most one
``role: "title"`` slot (matched to the slide's one scorable title
shape); ``role: "body"`` slots are matched, in the order written, to
the slide's scorable body shapes in their own ranked (best-first)
order -- the caller already judged this slide's real layout by eye, so
this positional match trusts that judgment rather than re-deriving it.
``text`` as a plain string injects one line (autofit/font-size-balanced
same as every other layout); as a list of strings, injects a bulleted
list into that one shape. Refuses loudly (never silently truncates or
drops) if more slots are declared than the slide has matching shapes.

``source_slide_idx`` (optional, int): explicitly picks which real
template slide (0-based, ``prs.slides[source_slide_idx]``) to clone
for this output slide, bypassing ``pick_source_slide_index``'s
geometry-based heuristic entirely for that slide. RECOMMENDED
workflow for quality-sensitive decks: render ``thumbnail_grid.py``
first, look at the labeled grid, then supply the exact index you
chose. The hard structural checks (a scorable title shape, and
whichever `layout` needs, enough body-candidate shapes) still run --
an explicit choice that genuinely can't support the layout is refused
loudly (exit 2), never silently guessed past. Omit this field to keep
the existing automatic heuristic exactly as before (default,
backward-compatible path, still a legitimate lower-effort fallback).

``image_path`` (optional, str, only meaningful for ``"image"``
layout): path to a real caller-supplied image file (resolved relative
to the content.json file's own directory if not absolute) that
replaces the cloned slide's actual picture -- not just leaving the
template's original photo untouched. Refused loudly if the file
doesn't exist or isn't a decodable raster image. A markdown
``![alt](path)`` image reference inside the slide's own ``markdown``
field is honored as a fallback source when ``image_path`` is omitted.

``generate_placeholder_image`` (optional, bool, only meaningful for
``"image"`` layout): when true, injects a deterministic, network-free
Pillow-generated placeholder (same engine as
``placeholder_images.generate_placeholder_png``) instead of requiring
a real file -- the no-network fallback path. Ignored if ``image_path``
is also set (the real file wins). If neither is set and the layout is
``"image"``, the template's own original picture is left untouched
exactly as before (full backward compatibility).

Exit codes: 0 = success, 2 = malformed/invalid input (refuse loudly,
never guess), 3 = template structurally unusable (e.g. zero slides).
"""
from __future__ import annotations

import argparse
import json
import sys
from pathlib import Path

from pptx import Presentation

import constants as C
import text_utils
import rich_text
import shape_roles
import slide_classifier
import slide_duplicator
import font_manager
import font_engine
import font_vietnamese
import font_size_balance
import image_injector

for _stream in (sys.stdout, sys.stderr):
    if hasattr(_stream, "reconfigure"):
        _stream.reconfigure(encoding="utf-8")

_DEFAULT_TITLE_PT = 28.0
_DEFAULT_BODY_PT = 16.0


class ComposerError(Exception):
    """Raised for any refuse-loudly condition (never guess, never silently drop)."""


def load_content_spec(path: Path) -> list[dict]:
    if not path.is_file():
        raise ComposerError(f"content file not found: {path}")
    try:
        data = json.loads(path.read_text(encoding="utf-8"))
    except json.JSONDecodeError as e:
        raise ComposerError(f"malformed JSON in {path}: {e}") from e
    if not isinstance(data, dict) or "slides" not in data:
        raise ComposerError("content JSON must be an object with a 'slides' array")
    slides = data["slides"]
    if not isinstance(slides, list) or not slides:
        raise ComposerError("content.slides must be a non-empty array")
    for i, item in enumerate(slides):
        if not isinstance(item, dict) or "layout" not in item:
            raise ComposerError(f"slides[{i}] must be an object with at least a 'layout' key")
        if item["layout"] not in C.LAYOUT_NAME_MAP:
            raise ComposerError(
                f"slides[{i}].layout={item['layout']!r} is not implemented. "
                f"Supported: {sorted(C.LAYOUT_NAME_MAP)}. See SKILL.md 'Known limitations' "
                f"for the layout types not yet built -- refusing rather than silently "
                f"downgrading to a layout you didn't ask for."
            )
        is_custom = C.LAYOUT_NAME_MAP[item["layout"]] == C.LAYOUT_CUSTOM
        if not is_custom and "markdown" not in item:
            raise ComposerError(f"slides[{i}] must have a 'markdown' key (all layouts except 'custom')")
        if "source_slide_idx" in item and item["source_slide_idx"] is not None:
            idx = item["source_slide_idx"]
            if not isinstance(idx, int) or isinstance(idx, bool) or idx < 0:
                raise ComposerError(
                    f"slides[{i}].source_slide_idx must be a non-negative integer slide "
                    f"index, got {idx!r}"
                )
        if "image_path" in item and item["image_path"] is not None and not isinstance(item["image_path"], str):
            raise ComposerError(f"slides[{i}].image_path must be a string path, got {item['image_path']!r}")
        if "generate_placeholder_image" in item and not isinstance(item["generate_placeholder_image"], bool):
            raise ComposerError(
                f"slides[{i}].generate_placeholder_image must be a boolean, "
                f"got {item['generate_placeholder_image']!r}"
            )
        if is_custom:
            # "custom" is the ad-hoc escape hatch for a real template
            # slide shape that doesn't match one of the 6 named layouts
            # (see SKILL.md "Custom layout") -- it has NO automatic
            # heuristic (there's nothing generic to search a pool for),
            # so an explicit source_slide_idx is mandatory, not optional.
            if item.get("source_slide_idx") is None:
                raise ComposerError(
                    f"slides[{i}]: layout='custom' requires an explicit source_slide_idx "
                    f"-- there is no automatic heuristic for an ad-hoc slot mapping. Render "
                    f"thumbnail_grid.py, look at the slide, then supply its index."
                )
            slots = item.get("slots")
            if not isinstance(slots, list) or not slots:
                raise ComposerError(f"slides[{i}]: layout='custom' requires a non-empty 'slots' array")
            title_slots = 0
            for j, slot in enumerate(slots):
                if not isinstance(slot, dict) or "role" not in slot or "text" not in slot:
                    raise ComposerError(f"slides[{i}].slots[{j}] must be an object with 'role' and 'text' keys")
                if slot["role"] not in ("title", "body"):
                    raise ComposerError(f"slides[{i}].slots[{j}].role must be 'title' or 'body', got {slot['role']!r}")
                text = slot["text"]
                if not (isinstance(text, str) or (isinstance(text, list) and all(isinstance(t, str) for t in text))):
                    raise ComposerError(
                        f"slides[{i}].slots[{j}].text must be a string (single line) or a list "
                        f"of strings (bullets), got {text!r}"
                    )
                if slot["role"] == "title":
                    title_slots += 1
            if title_slots > 1:
                raise ComposerError(
                    f"slides[{i}]: only one slots[].role='title' is allowed (a slide has one "
                    f"title shape), got {title_slots}"
                )
    return slides


# Minimum structurally-required body/content shapes for each layout,
# used to pick a STRUCTURALLY COMPATIBLE source slide from the template
# instead of blind round-robin. Bug found in self-test round 1 (see
# SKILL.md changelog): blind round-robin cloned a "standard" bullets
# request onto a two-column template slide (and vice versa), which (a)
# left the unused second placeholder's ORIGINAL DEMO TEXT visible in
# the output ("Right demo point" leaking into a real deck) and (b)
# silently dropped real caller content that had nowhere structurally
# compatible to go (a two-column right-side list vanishing entirely).
# Fixed two ways: this structural pre-filter, AND a post-injection
# sweep (see ``_clear_unused_text_shapes``) that blanks any leftover
# non-decorative text shape not actually used for this slide's content,
# so stale template demo text can never survive into real output even
# when a perfect structural match isn't available.
_MIN_BODY_REQUIRED: dict[str, int] = {
    C.LAYOUT_TITLE: 0,
    C.LAYOUT_STD: 1,
    C.LAYOUT_2COL: 2,
    C.LAYOUT_IMG: 1,
    C.LAYOUT_QUOTE: 0,
    C.LAYOUT_SECTION: 1,
}


def _title_and_body_counts(prs, slide_idx: int) -> tuple[bool, int]:
    title, bodies = shape_roles.find_title_and_body(
        prs.slides[slide_idx].shapes, prs.slide_width, prs.slide_height
    )
    return title is not None, len(bodies)


def _best_body_capacity(prs, slide_idx: int) -> int:
    """Largest capacity_chars among this slide's scored body candidates
    (0 if none)."""
    _title, bodies = shape_roles.find_title_and_body(
        prs.slides[slide_idx].shapes, prs.slide_width, prs.slide_height
    )
    if not bodies:
        return 0
    return max(b.get("capacity_chars", 0) for b in bodies)


# Slack factor applied to a candidate slide's best body capacity before
# comparing against the real content length: font-size balancing can
# still shrink text somewhat, so a slide doesn't need 100% headroom to
# be usable -- but it must clear at least this fraction, or the content
# is being crammed rather than fitted.
_CAPACITY_SLACK = 0.6


def _title_capacity(prs, slide_idx: int) -> int:
    title, _bodies = shape_roles.find_title_and_body(
        prs.slides[slide_idx].shapes, prs.slide_width, prs.slide_height
    )
    return title.get("capacity_chars", 0) if title else 0


def _fits(
    prs, slide_idx: int, min_body: int, min_capacity_chars: int = 0,
    min_title_capacity_chars: int = 0, max_unused_body_shapes: int | None = None,
) -> bool:
    """A candidate source slide is usable only when it has a scorable
    title shape, enough body candidates, AND (when real content length
    is known) at least one body candidate with real character capacity
    for that content.

    Two real bugs found this way:
    1. (self-test round 3, "Elegant Black & White Style" template) its
       closing slide is a pure full-bleed picture (BLANK layout, a
       single Picture shape, zero text frames) -- min_body=0 for a
       "quote" layout let it through the old filter (body count only),
       so the quote's whole title text silently had nowhere to go and
       vanished. Fixed with the has_title check below.
    2. (hands-on visual review, rendered to JPG and inspected) a
       template's 6-item icon-grid slide has 6 tiny CONTENT-typed
       caption placeholders that satisfy any min_body>=1 count check,
       so a real multi-bullet paragraph kept getting crammed into one
       tiny caption box while 6 unrelated icons dominated the slide
       visually. Counting body shapes was never enough -- their actual
       SIZE has to be checked against the real content about to be
       injected. Fixed with the capacity check below.
    """
    has_title, body_count = _title_and_body_counts(prs, slide_idx)
    if not (has_title and body_count >= min_body):
        return False
    if min_title_capacity_chars > 0 and _title_capacity(prs, slide_idx) < min_title_capacity_chars * _CAPACITY_SLACK:
        return False
    if max_unused_body_shapes is not None and body_count > max_unused_body_shapes:
        # Real bug found in a third visual review: a QUOTE/TITLE layout
        # only injects into the title shape, so a candidate slide's body
        # capacity is never actually checked for those layouts -- but a
        # 6-icon-grid slide still passes every other check (it has a
        # normal-sized title placeholder, separate from its 6 tiny icon
        # captions) and gets picked, leaving 6 unrelated icons visually
        # dominating an otherwise-clean quote slide. A title/quote slide
        # is expected to be visually minimal; a slide with many small
        # body/caption shapes is a structural signal of an icon-grid or
        # stat-card design, not a clean statement slide -- reject it
        # here regardless of whether the title text itself would fit.
        return False
    if min_capacity_chars <= 0:
        return True
    return _best_body_capacity(prs, slide_idx) >= min_capacity_chars * _CAPACITY_SLACK


# Absolute floor for any layout that injects real prose into a body
# shape, independent of how short the caller's actual text happens to
# be. Without this, a short section subtitle (e.g. "Why Scriptorium
# exists", ~23 chars) still clears a tiny icon-grid caption box's
# capacity via the relative _CAPACITY_SLACK check, but the box remains
# the wrong SHAPE for prose -- surrounded by 6 unrelated decorative
# icons that stay visually dominant regardless of how little text goes
# in. A real prose box, even for one short line, should be sized like
# a sentence, not a numeral badge. Found via a second hands-on visual
# review after the first capacity fix only helped the longer slides.
_MIN_PROSE_CAPACITY_CHARS = 120

# A normal single-message template slide (title+bullets, two-column,
# caption+image, section divider) has a small, coherent number of body
# candidates -- typically 1-4. A slide scoring far more than that is a
# strong structural signal of an infographic/map/diagram slide (many
# small per-region/per-item labels scattered over a decorative
# graphic), not a plain content slide -- wrong shape for prose
# injection regardless of any single candidate's capacity. Found via a
# fourth hands-on visual review: a real Slidesgo template slide with a
# decorative world-map graphic and 19 scored body candidates (small
# per-country label boxes) got picked for a plain 6-bullet "standard"
# slide, producing a map irrelevant to the content with bullets jammed
# into one of the tiny label boxes.
_MAX_BODY_SHAPES_FOR_PROSE = 8


def _required_body_capacity(layout_type: str, content) -> int:
    """Character length of the largest body block this layout will try
    to inject, used to reject a structurally-plausible-but-too-small
    source slide before it's picked. 0 means "no body content to size
    against" (title/quote layouts check title capacity separately, not
    via this path)."""
    if layout_type == C.LAYOUT_STD:
        raw = len("\n".join(content.bullets))
    elif layout_type == C.LAYOUT_2COL:
        left = "\n".join([content.left_header, *content.left_column]) if content.left_header else "\n".join(content.left_column)
        right = "\n".join([content.right_header, *content.right_column]) if content.right_header else "\n".join(content.right_column)
        raw = max(len(left), len(right))
    elif layout_type == C.LAYOUT_IMG:
        raw = len(content.caption)
    elif layout_type == C.LAYOUT_SECTION:
        raw = len(content.subtitle or " ".join(content.bullets))
    elif layout_type == C.LAYOUT_TITLE:
        raw = len(content.subtitle)
    else:
        return 0
    if raw <= 0:
        return 0
    return max(raw, _MIN_PROSE_CAPACITY_CHARS)


def _required_title_capacity(layout_type: str, content) -> int:
    """Quote content is injected into the TITLE shape, not a body shape
    (see compile_deck's LAYOUT_QUOTE branch) -- a long standalone quote
    needs a title shape sized for a real sentence, not just any slide
    with a scorable title. Other layouts' titles are short headings by
    convention and don't need this extra check."""
    if layout_type == C.LAYOUT_QUOTE and content.title:
        return max(len(content.title), _MIN_PROSE_CAPACITY_CHARS)
    return 0


# Real bug found via a real 20-slide showcase-deck dogfooding round done
# in a SEPARATE session (v0.8.0, 2026-08-03, test/Aug032026/): an agent
# that HAD looked at a rendered thumbnail grid still explicitly chose a
# real "3 parallel ideas" template slide (3 numbered-badge + heading +
# description card clusters, 6 real scorable body shapes total) as the
# source for a plain 2-slot `two_column` request. `_hard_structural_fit`
# only checked `body_count >= min_body` (2), which a 6-body-shape slide
# trivially clears -- but `find_title_and_body`'s ranking picks its N
# best-SCORING body shapes in isolation, with no notion that 2 of those
# 6 shapes are a matched header+description PAIR belonging to one card
# and the other 4 belong to two entirely different cards. The result,
# confirmed by rendering: content landed in 2 shapes that don't even
# belong to the same card, while the other 2 cards' badges ("2", "3")
# were left with headings/descriptions completely blank -- a much worse
# defect than a stale-text leak, since real content silently vanished
# from 2 of the deck's visual "cards" with no error anywhere in the
# compile report. A rendered thumbnail is too small to show this at a
# glance -- a 3-card cluster and a genuine 6-item bullet list can look
# identical in a scaled-down grid cell. Fixed with the same "shape-count
# is a structural signal" principle proven in bugs #6-9: a layout that
# only needs `min_body` slots being offered several times that many real
# body shapes is itself evidence of a multi-card cluster the layout's
# flat slot model cannot address correctly, independent of what the
# thumbnail looked like. `_MAX_BODY_SLACK` caps how far above `min_body`
# an explicitly-chosen slide may go before this refuses loudly and
# names `custom` (one slot per real shape) as the correct tool instead.
_MAX_BODY_SLACK: dict[str, int] = {
    C.LAYOUT_TITLE: 2,
    C.LAYOUT_STD: _MAX_BODY_SHAPES_FOR_PROSE,
    C.LAYOUT_2COL: 3,
    C.LAYOUT_IMG: 4,
    C.LAYOUT_QUOTE: 2,
    C.LAYOUT_SECTION: 4,
}


def _hard_structural_fit(
    prs, slide_idx: int, min_body: int, *, layout_type: str | None = None,
) -> tuple[bool, str]:
    """The minimum, non-negotiable structural bar an EXPLICITLY chosen
    ``source_slide_idx`` must clear: a scorable title shape, at least
    as many scorable body candidates as the requested layout needs,
    and (see ``_MAX_BODY_SLACK`` above) not drastically more than that
    either -- a caller who already looked at a rendered thumbnail grid
    has judged visual fit directly, but cannot see from a small
    thumbnail whether a slide's body shapes form one flat pool or
    several unrelated card clusters, so that specific check still runs
    even on an explicit choice. ``layout_type`` is optional only for
    ``LAYOUT_CUSTOM`` callers, which pass their own exact declared slot
    count as ``min_body`` and are exempt from the slack cap -- custom
    already addresses every real shape it uses by an explicit slot, so
    there is no flat-pool-vs-cluster ambiguity to guard against there.
    Returns (ok, reason) so the caller can refuse loudly with a
    concrete explanation rather than a bare boolean."""
    has_title, body_count = _title_and_body_counts(prs, slide_idx)
    if not has_title:
        return False, "no scorable title shape found on this slide"
    if body_count < min_body:
        return False, f"only {body_count} scorable body shape(s) found, this layout needs >= {min_body}"
    if layout_type is not None and layout_type != C.LAYOUT_CUSTOM:
        max_body = min_body + _MAX_BODY_SLACK.get(layout_type, _MAX_BODY_SHAPES_FOR_PROSE)
        if body_count > max_body:
            display_name = next(
                (k for k, v in C.LAYOUT_NAME_MAP.items() if v == layout_type), layout_type
            )
            return False, (
                f"{body_count} scorable body shape(s) found, but layout {display_name!r} only "
                f"uses {min_body} of them -- this many real shapes on one slide usually means "
                f"several parallel content 'cards' (e.g. a 3-column icon/stat layout), which a "
                f"flat {display_name!r} slot model will fill unevenly, leaving some cards blank. "
                f"Use the 'custom' layout with one slot per real shape instead, or choose a "
                f"source slide with fewer body shapes."
            )
    return True, ""


def pick_source_slide_index(
    prs,
    layout_type: str,
    position: int,
    total: int,
    classification: dict[str, list[int]],
    original_count: int,
    layout_round_robin: dict[str, int],
    required_capacity_chars: int = 0,
    required_title_capacity_chars: int = 0,
) -> tuple[int, bool]:
    """Pick a structurally-compatible template slide to clone for this
    output slide. Returns (source_slide_idx, structural_match)."""
    title_pool = classification.get("title", [])
    content_pool = classification.get("content", [])
    closing_pool = classification.get("closing", [])
    all_pool = title_pool + content_pool + closing_pool
    if not all_pool:
        all_pool = list(range(original_count))

    min_body = _MIN_BODY_REQUIRED.get(layout_type, 0)

    if layout_type == C.LAYOUT_TITLE:
        candidates = [
            i for i in (title_pool or content_pool or all_pool)
            if _fits(prs, i, 0, max_unused_body_shapes=2)
        ]
        if not candidates:
            candidates = [i for i in (title_pool or content_pool or all_pool) if _fits(prs, i, 0)]
        if candidates:
            return candidates[0], True
        # No title-pool slide has a scorable title shape (rare, but a
        # picture-only cover slide is possible) -- fall through to the
        # general search below instead of guaranteeing a content-loss.
        title_pool = []

    # Only reuse a "closing" slide (typically minimal/single-message,
    # e.g. "Thank you") for layouts that themselves need at most a
    # title -- and only when that slide actually HAS a title-capable
    # shape (see _fits docstring for the real bug this guards against).
    if position == total - 1 and closing_pool and min_body == 0:
        closing_candidates = [i for i in closing_pool if _fits(prs, i, min_body)]
        if closing_candidates:
            return closing_candidates[0], True

    search_pool = content_pool or all_pool

    # Capacity-aware search first: reject slides whose best body shape
    # is too small for the real content about to be injected (see
    # _fits docstring, bug 2). Only falls back to the capacity-blind
    # search when NO slide in the template clears the bar -- a
    # template that is uniformly tight on space shouldn't refuse to
    # compile, but a template that genuinely HAS a spacious slide
    # should never lose to a tiny one via round-robin luck.
    # Quote/title-only layouts never touch body shapes, so a slide with
    # many small body/caption shapes (an icon-grid/stat-card design)
    # would otherwise pass on title capacity alone while leaving those
    # unrelated shapes visually dominant. Every layout now shares the
    # same per-layout cap `_hard_structural_fit` enforces on an explicit
    # choice (see `_MAX_BODY_SLACK` and its v0.8.0 docstring) -- the
    # automatic path was previously looser (a flat 8-shape cap for every
    # non-quote/title layout), which is exactly wide enough to have let
    # the real card-cluster mismatch bug through here too, not just via
    # an explicit source_slide_idx.
    max_unused = min_body + _MAX_BODY_SLACK.get(layout_type, _MAX_BODY_SHAPES_FOR_PROSE)

    if required_capacity_chars > 0 or required_title_capacity_chars > 0 or max_unused is not None:
        capacity_matching = [
            i for i in search_pool
            if _fits(prs, i, min_body, required_capacity_chars, required_title_capacity_chars, max_unused)
        ]
        if capacity_matching:
            counter = layout_round_robin.get(layout_type, 0)
            layout_round_robin[layout_type] = counter + 1
            return capacity_matching[counter % len(capacity_matching)], True

    matching = [i for i in search_pool if _fits(prs, i, min_body)]
    if not matching:
        matching = [i for i in all_pool if _fits(prs, i, min_body)]

    if matching:
        counter = layout_round_robin.get(layout_type, 0)
        layout_round_robin[layout_type] = counter + 1
        return matching[counter % len(matching)], True

    # No structurally-compatible slide anywhere in the template: fall
    # back to the richest available slide rather than crashing the
    # whole compile over one slide, but report the mismatch (never
    # silently pretend the content fit).
    best = max(search_pool, key=lambda i: _title_and_body_counts(prs, i)[1])
    return best, False


def _clear_shape_text(shape) -> None:
    """Blank a shape's text down to a single empty paragraph, XML-level
    (matches the wipe pattern used when python-pptx add_slide() first
    injects placeholder demo text from the layout)."""
    tf = shape.text_frame
    txBody = tf._txBody
    for p in list(txBody.findall(f"{C.NS_A}p")):
        txBody.remove(p)
    from lxml import etree
    empty_p = etree.SubElement(txBody, f"{C.NS_A}p")
    etree.SubElement(empty_p, f"{C.NS_A}endParaRPr")


def _clear_unused_text_shapes(all_infos: list[dict], used_shape_ids: set[int]) -> int:
    """Blank any non-decorative, non-watermark text shape on the cloned
    slide that was NOT assigned a role for this slide's content, so
    the template's own original demo/instruction text never leaks into
    real output. Returns the number of shapes cleared."""
    cleared = 0
    for info in all_infos:
        is_vendor_caption = shape_roles.is_vendor_instructional_caption(info["text"])
        if (info["is_decorative"] or info["is_watermark"]) and not is_vendor_caption:
            continue
        if id(info["shape"]._element) in used_shape_ids:
            continue
        if info["text"]:
            _clear_shape_text(info["shape"])
            cleared += 1
    return cleared


def _resolve_font_name(
    requested: str | None,
    theme_fallback: str,
    vietnamese: bool,
    role_label: str,
    log: list[dict],
) -> str | None:
    """Resolve the font name to actually WRITE onto runs.

    Returns None when no explicit override is warranted (leave the
    template's own inherited font alone -- the default, most-common
    path). Returns a concrete font name when the caller requested a
    swap, or when --vietnamese forced a substitution because the
    effective font lacks confirmed Vietnamese support.
    """
    effective = requested or theme_fallback
    if not effective:
        return requested  # nothing to check against; honor caller's raw request or None

    if not vietnamese:
        return requested  # only override runs when caller explicitly requested a font swap

    check = font_vietnamese.check_font_vietnamese_safety(effective)
    if check.status == "unknown":
        raise ComposerError(
            f"--vietnamese requested but the {role_label} font '{effective}' has unknown "
            f"Vietnamese glyph coverage: {check.reason}"
        )
    resolved = check.substitute if check.status == "substitute" else check.font_name
    log.append({
        "role": role_label, "requested_or_theme": effective,
        "status": check.status, "resolved": resolved,
    })
    return resolved


def _inject_single_line(
    shape, text: str, *, font_name: str | None, font_size_pt: float, bold: bool = False,
) -> None:
    tf = shape.text_frame
    while len(tf.paragraphs) > 1:
        tf._txBody.remove(tf.paragraphs[-1]._p)
    p = tf.paragraphs[0]
    runs = rich_text.parse_inline_formatting(text_utils.clean_markdown_formatting(text))
    rich_text.apply_runs_to_paragraph(
        p, runs, base_font_name=font_name, base_font_size_pt=font_size_pt, base_bold=bold,
    )
    text_utils.enable_autofit(tf)
    text_utils.set_text_wrap(tf)


def _inject_bullets(
    shape, bullets: list[str], *, font_name: str | None, font_size_pt: float,
    heading: str | None = None,
) -> None:
    """Inject a bullet list into `shape`'s text frame, optionally
    preceded by a bold heading paragraph in the same shape.

    `heading` closes a real bug found via a cold-agent usability test
    (2026-08-02): text_utils._parse_two_column_slide already parses a
    two_column slide's "### Left"/"### Right" markdown into
    ParsedSlideContent.left_header/right_header, but LAYOUT_2COL's
    injection branch never read either field -- the sub-heading text
    was silently discarded with no error, no warning, and not counted
    in stale_text_shapes_cleared. The tester (a fresh agent with no
    prior context on this skill, given a real client brief) found this
    the hard way and worked around it by folding the heading into the
    first bullet line; that workaround shouldn't be necessary. This
    shape only has ONE text frame per column (no separate heading
    placeholder to inject into), so the fix is to inject the heading
    as a bold first paragraph inside the same shape, ahead of the
    bullets -- not a second shape, since none exists for it.
    """
    tf = shape.text_frame
    while len(tf.paragraphs) > 1:
        tf._txBody.remove(tf.paragraphs[-1]._p)
    if not bullets:
        bullets = [""]
    first = True
    if heading and heading.strip():
        p = tf.paragraphs[0]
        clean = text_utils.clean_markdown_formatting(heading)
        runs = rich_text.parse_inline_formatting(clean)
        rich_text.apply_runs_to_paragraph(p, runs, base_font_name=font_name, base_font_size_pt=font_size_pt, base_bold=True)
        first = False
    for bullet in bullets:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        clean = text_utils.clean_markdown_formatting(bullet)
        runs = rich_text.parse_inline_formatting(clean)
        rich_text.apply_runs_to_paragraph(p, runs, base_font_name=font_name, base_font_size_pt=font_size_pt)
    text_utils.enable_autofit(tf)
    text_utils.set_text_wrap(tf)


def compile_deck(
    template_path: Path,
    content_path: Path,
    output_path: Path,
    heading_font: str | None,
    body_font: str | None,
    vietnamese: bool,
) -> dict:
    slides_spec = load_content_spec(content_path)

    if not template_path.is_file():
        raise ComposerError(f"template file not found: {template_path}")
    prs = Presentation(str(template_path))

    original_count = len(prs.slides)
    if original_count == 0:
        raise ComposerError("template has zero slides -- nothing to clone from")

    classification = slide_classifier.classify_template_slides(prs, original_count)

    theme = font_manager.FontManager(prs).get_theme_fonts()
    font_log: list[dict] = []
    # Only pre-resolve a GLOBAL font when the caller explicitly requested
    # a font swap (--heading-font/--body-font) -- that's a real request
    # to validate and apply upfront, unlike the plain --vietnamese case,
    # which is now handled per-shape below (_font_for) using each
    # shape's own real effective font (see font_manager.resolve_
    # effective_font's docstring, v0.6.0) rather than blindly checking
    # only the theme's fontScheme for every shape in the deck.
    resolved_heading_font: str | None = None
    resolved_body_font: str | None = None
    if heading_font or body_font:
        resolved_heading_font = _resolve_font_name(heading_font, theme.major_font, vietnamese, "heading", font_log)
        resolved_body_font = _resolve_font_name(body_font, theme.minor_font, vietnamese, "body", font_log)
        applied_heading = resolved_heading_font or theme.major_font
        applied_body = resolved_body_font or theme.minor_font
        if applied_heading and applied_body:
            font_engine.set_theme_fonts(prs, applied_heading, applied_body)

    total = len(slides_spec)
    slide_plan: list[dict] = []
    size_decisions: list[dict] = []
    layout_round_robin: dict[str, int] = {}

    for position, item in enumerate(slides_spec):
        layout_type = C.LAYOUT_NAME_MAP[item["layout"]]

        # Parsed BEFORE source-slide selection so pick_source_slide_index
        # can reject a structurally-plausible-but-too-small candidate
        # (see _fits's bug-2 docstring) -- selection needs to know the
        # real content length, not just the layout type. "custom" has no
        # markdown at all (content lives in item["slots"] instead), and
        # no automatic selection either (source_slide_idx is mandatory
        # for it, enforced in load_content_spec), so this whole
        # content-parsing/capacity-estimation step is meaningless for it.
        if layout_type == C.LAYOUT_CUSTOM:
            content = None
            required_capacity = 0
            required_title_capacity = 0
        else:
            content = text_utils.parse_slide_content(layout_type, item["markdown"])
            required_capacity = _required_body_capacity(layout_type, content)
            required_title_capacity = _required_title_capacity(layout_type, content)

        explicit_idx = item.get("source_slide_idx")
        if explicit_idx is not None:
            # RECOMMENDED workflow for quality-sensitive decks (see
            # thumbnail_grid.py + SKILL.md): the caller already looked
            # at a rendered thumbnail grid and deliberately chose this
            # exact slide -- skip the geometry heuristic entirely, but
            # still refuse loudly if the choice is structurally
            # impossible for the requested layout (never silently
            # guess past a real mismatch).
            if explicit_idx >= original_count:
                raise ComposerError(
                    f"slides[{position}].source_slide_idx={explicit_idx} is out of range "
                    f"-- template has {original_count} slide(s) (valid: 0..{original_count - 1})"
                )
            if layout_type == C.LAYOUT_CUSTOM:
                # Real required body-slot count comes from THIS slide's
                # own declared slots, not a static per-layout table --
                # every "custom" slide can ask for a different number.
                min_body_explicit = sum(1 for s in item["slots"] if s["role"] == "body")
            else:
                min_body_explicit = _MIN_BODY_REQUIRED.get(layout_type, 0)
            ok, reason = _hard_structural_fit(prs, explicit_idx, min_body_explicit, layout_type=layout_type)
            if not ok:
                raise ComposerError(
                    f"slides[{position}]: explicit source_slide_idx={explicit_idx} does not "
                    f"structurally support layout {item['layout']!r}: {reason}"
                )
            source_idx, structural_match = explicit_idx, True
        else:
            source_idx, structural_match = pick_source_slide_index(
                prs, layout_type, position, total, classification, original_count, layout_round_robin,
                required_capacity_chars=required_capacity,
                required_title_capacity_chars=required_title_capacity,
            )
        source_slide = prs.slides[source_idx]
        new_slide = slide_duplicator.duplicate_slide(prs, source_slide)

        # Captured once, before any injection, so (a) role scoring sees
        # the template's real original text/geometry and (b) the later
        # stale-text cleanup sweep knows what every shape looked like
        # before this slide's content touched it.
        all_text_infos = shape_roles.collect_text_shapes(new_slide.shapes, prs.slide_width, prs.slide_height)
        title_info, body_infos = shape_roles.find_title_and_body(
            new_slide.shapes, prs.slide_width, prs.slide_height
        )
        used_shape_ids: set[int] = set()
        if title_info is not None:
            used_shape_ids.add(id(title_info["shape"]._element))

        record = {
            "position": position, "layout": item["layout"], "source_slide_idx": source_idx,
            "structural_match": structural_match,
            "source_slide_idx_explicit": explicit_idx is not None,
            "title_shape_found": title_info is not None, "body_shapes_found": len(body_infos),
        }
        slide_plan.append(record)

        def _size_for(info, new_text: str, role: str) -> float:
            base_pt = info["font_size_pt"] or (_DEFAULT_TITLE_PT if role == "title" else _DEFAULT_BODY_PT)
            decision = font_size_balance.compute_balanced_font_size(
                original_char_count=info["char_count"],
                new_char_count=len(new_text),
                original_pt=base_pt,
                role=role,
            )
            if decision.applied:
                size_decisions.append({
                    "position": position, "role": role,
                    "original_chars": decision.original_char_count,
                    "new_chars": decision.new_char_count,
                    "ratio": round(decision.ratio, 2),
                    "steps": decision.steps,
                    "original_pt": decision.original_pt,
                    "balanced_pt": decision.balanced_pt,
                })
            return decision.balanced_pt

        def _font_for(shape, role: str) -> str | None:
            """Per-SHAPE font resolution (v0.6.0), replacing a single
            deck-wide theme-only decision. See font_manager.
            resolve_effective_font's docstring for the real bug this
            fixes (found via a cold-agent usability test): a template
            can declare its real display font at the slide-layout or
            master level, invisible to a theme-only check, and forcing
            the theme font onto every run discarded that real
            typography even when the theme font itself was 'safe'.
            """
            explicit = heading_font if role == "title" else body_font
            if explicit:
                # An explicit --heading-font/--body-font request is a
                # real global override -- applies uniformly, already
                # resolved (and validated) once before this loop.
                return resolved_heading_font if role == "title" else resolved_body_font
            if not vietnamese:
                return None  # no override requested at all -- leave inherited font alone
            shape_font = font_manager.resolve_effective_font(shape, new_slide)
            theme_font = theme.major_font if role == "title" else theme.minor_font
            effective = shape_font or theme_font
            if not effective:
                return None
            check = font_vietnamese.check_font_vietnamese_safety(effective)
            if check.status == "unknown":
                raise ComposerError(
                    f"slides[{position}]: --vietnamese requested but the {role} shape's "
                    f"effective font '{effective}' has unknown Vietnamese glyph coverage: "
                    f"{check.reason}"
                )
            resolved = check.substitute if check.status == "substitute" else check.font_name
            font_log.append({
                "position": position, "role": role, "requested_or_theme": effective,
                "status": check.status, "resolved": resolved,
                "source": "shape/layout/master" if shape_font else "theme",
            })
            return resolved

        if layout_type == C.LAYOUT_TITLE:
            if title_info is not None:
                pt = _size_for(title_info, content.title, "title")
                _inject_single_line(title_info["shape"], content.title, font_name=_font_for(title_info["shape"], "title"), font_size_pt=pt)
            if content.subtitle and body_infos:
                sub_info = body_infos[0]
                pt = _size_for(sub_info, content.subtitle, "body")
                _inject_single_line(sub_info["shape"], content.subtitle, font_name=_font_for(sub_info["shape"], "body"), font_size_pt=pt)
                used_shape_ids.add(id(sub_info["shape"]._element))

        elif layout_type == C.LAYOUT_STD:
            if title_info is not None:
                pt = _size_for(title_info, content.title, "title")
                _inject_single_line(title_info["shape"], content.title, font_name=_font_for(title_info["shape"], "title"), font_size_pt=pt)
            if body_infos:
                body_info = body_infos[0]
                joined = "\n".join(content.bullets)
                pt = _size_for(body_info, joined, "body")
                _inject_bullets(body_info["shape"], content.bullets, font_name=_font_for(body_info["shape"], "body"), font_size_pt=pt)
                used_shape_ids.add(id(body_info["shape"]._element))

        elif layout_type == C.LAYOUT_2COL:
            if title_info is not None:
                pt = _size_for(title_info, content.title, "title")
                _inject_single_line(title_info["shape"], content.title, font_name=_font_for(title_info["shape"], "title"), font_size_pt=pt)
            cols = sorted(body_infos[:2], key=lambda i: i["left_pct"])
            if len(cols) >= 1 and content.left_column:
                sized_text = "\n".join([content.left_header, *content.left_column]) if content.left_header else "\n".join(content.left_column)
                pt = _size_for(cols[0], sized_text, "body")
                _inject_bullets(cols[0]["shape"], content.left_column, font_name=_font_for(cols[0]["shape"], "body"), font_size_pt=pt, heading=content.left_header)
                used_shape_ids.add(id(cols[0]["shape"]._element))
            if len(cols) >= 2 and content.right_column:
                sized_text = "\n".join([content.right_header, *content.right_column]) if content.right_header else "\n".join(content.right_column)
                pt = _size_for(cols[1], sized_text, "body")
                _inject_bullets(cols[1]["shape"], content.right_column, font_name=_font_for(cols[1]["shape"], "body"), font_size_pt=pt, heading=content.right_header)
                used_shape_ids.add(id(cols[1]["shape"]._element))

        elif layout_type == C.LAYOUT_IMG:
            if title_info is not None:
                pt = _size_for(title_info, content.title, "title")
                _inject_single_line(title_info["shape"], content.title, font_name=_font_for(title_info["shape"], "title"), font_size_pt=pt)
            if body_infos and content.caption:
                cap_info = body_infos[0]
                pt = _size_for(cap_info, content.caption, "body")
                _inject_single_line(cap_info["shape"], content.caption, font_name=_font_for(cap_info["shape"], "body"), font_size_pt=pt)
                used_shape_ids.add(id(cap_info["shape"]._element))

            # Real picture injection (see image_injector.py). Only
            # attempted when the caller actually asked for it --
            # image_path / generate_placeholder_image / a markdown
            # ![alt](path) reference -- otherwise the template's own
            # original picture is left exactly as cloned, unchanged
            # from this skill's pre-existing behavior (full backward
            # compatibility with every earlier content.json).
            wants_placeholder = bool(item.get("generate_placeholder_image", False))
            image_path_raw = item.get("image_path") or content.image_ref
            if wants_placeholder or image_path_raw:
                resolved_image_path = None
                if image_path_raw and not wants_placeholder:
                    candidate = Path(image_path_raw)
                    if not candidate.is_absolute():
                        candidate = content_path.parent / candidate
                    if not candidate.is_file():
                        raise ComposerError(
                            f"slides[{position}]: image_path/image_ref {image_path_raw!r} "
                            f"not found (resolved to {candidate})"
                        )
                    resolved_image_path = candidate
                seed = content.caption or content.title or f"slide-{position}"
                try:
                    image_report = image_injector.inject_slide_image(
                        new_slide, prs.slide_width, prs.slide_height,
                        image_path=resolved_image_path, placeholder_seed=seed,
                    )
                except image_injector.ImageInjectorError as e:
                    raise ComposerError(f"slides[{position}]: image injection failed: {e}") from e
                if image_report is None:
                    raise ComposerError(
                        f"slides[{position}]: image injection requested but the source "
                        f"slide (idx={source_idx}) has no top-level picture shape to "
                        f"inject into -- refusing rather than silently skipping"
                    )
                record["image_injection"] = image_report

        elif layout_type == C.LAYOUT_QUOTE:
            if title_info is not None:
                pt = _size_for(title_info, content.title, "title")
                _inject_single_line(title_info["shape"], content.title, font_name=_font_for(title_info["shape"], "title"), font_size_pt=pt)

        elif layout_type == C.LAYOUT_SECTION:
            if title_info is not None:
                pt = _size_for(title_info, content.title, "title")
                _inject_single_line(title_info["shape"], content.title, font_name=_font_for(title_info["shape"], "title"), font_size_pt=pt)
            # Real bug found in self-test round 1: markdown with a single
            # "## Heading\nplain text line" (no second heading) parses
            # via text_utils' standard-slide path, which routes an
            # unheaded body line into ``bullets``, not ``subtitle`` (only
            # an explicit H1+H2 pair populates ``subtitle`` -- see
            # _parse_standard_slide). LAYOUT_SECTION only checked
            # ``content.subtitle``, so that line silently vanished from
            # the output. Fixed: fall back to the bullets when no
            # explicit subtitle was parsed.
            subtitle_text = content.subtitle or " ".join(content.bullets)
            if subtitle_text and body_infos:
                sub_info = body_infos[0]
                pt = _size_for(sub_info, subtitle_text, "body")
                _inject_single_line(sub_info["shape"], subtitle_text, font_name=_font_for(sub_info["shape"], "body"), font_size_pt=pt)
                used_shape_ids.add(id(sub_info["shape"]._element))

        elif layout_type == C.LAYOUT_CUSTOM:
            # Ad-hoc slot mapping for a real template slide shape that
            # doesn't match one of the 6 named layouts (see SKILL.md
            # "Custom layout"). load_content_spec already validated
            # slots is a non-empty list with at most one role="title";
            # _hard_structural_fit already confirmed enough body shapes
            # exist. Positional match: the declared title slot (if any)
            # -> title_info; body slots, in the order the caller wrote
            # them, -> body_infos in shape_roles' own ranked order (best
            # candidate first) -- the caller chose that ranking implicitly
            # by looking at the rendered thumbnail and picking this slide.
            # Real friction found via a cold-agent usability test
            # (2026-08-02): a caller has no way to know which physical
            # shape a given body slot's RANKED position corresponds to
            # before compiling -- the thumbnail grid shows the whole
            # slide, not each candidate shape's individual rank, so a
            # tester matched slots to a team-photo slide by visual
            # left-to-right reading and got a male founder's bio
            # swapped under a female stock photo, only caught by
            # rendering and looking. slot_mapping (added to `record`
            # below) reports each slot's actual target shape position
            # (left_pct/top_pct/area_pct, from the same shape_roles
            # metadata already computed) so a caller can cross-check
            # against a rendered thumbnail WITHOUT a second compile.
            slot_mapping: list[dict] = []
            body_slot_idx = 0
            for slot in item["slots"]:
                text = slot["text"]
                is_bullets = isinstance(text, list)
                joined = "\n".join(text) if is_bullets else text
                preview = (text[0] if is_bullets and text else joined)[:60]
                if slot["role"] == "title":
                    if title_info is None:
                        raise ComposerError(
                            f"slides[{position}]: slots declares role='title' but source "
                            f"slide {source_idx} has no scorable title shape"
                        )
                    pt = _size_for(title_info, joined, "title")
                    if is_bullets:
                        _inject_bullets(title_info["shape"], text, font_name=_font_for(title_info["shape"], "title"), font_size_pt=pt)
                    else:
                        _inject_single_line(title_info["shape"], text, font_name=_font_for(title_info["shape"], "title"), font_size_pt=pt)
                    slot_mapping.append({
                        "role": "title", "text_preview": preview,
                        "target_left_pct": round(title_info["left_pct"], 3),
                        "target_top_pct": round(title_info["top_pct"], 3),
                        "target_area_pct": round(title_info["area_pct"], 3),
                    })
                else:
                    if body_slot_idx >= len(body_infos):
                        raise ComposerError(
                            f"slides[{position}]: slots declares more role='body' entries than "
                            f"source slide {source_idx} has scorable body shapes "
                            f"({len(body_infos)} available)"
                        )
                    info = body_infos[body_slot_idx]
                    body_slot_idx += 1
                    slot_mapping.append({
                        "role": "body", "text_preview": preview,
                        "target_left_pct": round(info["left_pct"], 3),
                        "target_top_pct": round(info["top_pct"], 3),
                        "target_area_pct": round(info["area_pct"], 3),
                    })
                    pt = _size_for(info, joined, "body")
                    if is_bullets:
                        _inject_bullets(info["shape"], text, font_name=_font_for(info["shape"], "body"), font_size_pt=pt)
                    else:
                        _inject_single_line(info["shape"], text, font_name=_font_for(info["shape"], "body"), font_size_pt=pt)
                    used_shape_ids.add(id(info["shape"]._element))
            record["slot_mapping"] = slot_mapping

        record["stale_text_shapes_cleared"] = _clear_unused_text_shapes(all_text_infos, used_shape_ids)

    slide_duplicator.strip_template_slides(prs, original_count)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))

    # Real finding from the same v0.8.0 dogfooding round: a real deck
    # reused the same source slide (e.g. one section-divider design for
    # 3 of 5 different parts, one bullet-slide design for 4 unrelated
    # content slides) with no signal anywhere that it had happened --
    # visually monotonous at best, and actively misleading when the
    # reused source slide carries its own baked-in decorative ordinal
    # (a "01" badge on a slide reused for "Part V") that this skill
    # cannot rewrite (see SKILL.md's "Known limitations"). This can't
    # be silently auto-fixed -- reusing a source slide is sometimes the
    # genuinely right choice when content truly repeats a shape -- but
    # the caller (or a reviewing agent) should see it without having to
    # cross-reference slide_plan by hand.
    source_idx_counts: dict[int, list[int]] = {}
    for rec in slide_plan:
        source_idx_counts.setdefault(rec["source_slide_idx"], []).append(rec["position"])
    repeated_source_slides = [
        {"source_slide_idx": idx, "positions": positions, "count": len(positions)}
        for idx, positions in sorted(source_idx_counts.items())
        if len(positions) > 1
    ]

    return {
        "output": str(output_path),
        "slides_generated": len(slides_spec),
        "template_classification": classification,
        "theme_fonts": {"major": theme.major_font, "minor": theme.minor_font},
        "font_resolution_log": font_log,
        "slide_plan": slide_plan,
        "font_size_balancing_applied": size_decisions,
        "repeated_source_slides": repeated_source_slides,
    }


def main() -> int:
    parser = argparse.ArgumentParser(description="Clone a real .pptx template and inject content.")
    parser.add_argument("template", type=Path, help="Caller-supplied .pptx template file")
    parser.add_argument("content", type=Path, help="content.json (see module docstring for shape)")
    parser.add_argument("output", type=Path, help="Output .pptx path")
    parser.add_argument("--heading-font", default=None, help="Override template's heading/title font")
    parser.add_argument("--body-font", default=None, help="Override template's body font")
    parser.add_argument(
        "--vietnamese", action="store_true",
        help="Require Vietnamese-diacritic-safe fonts; substitutes or refuses loudly if unresolvable",
    )
    args = parser.parse_args()

    try:
        result = compile_deck(
            args.template, args.content, args.output,
            args.heading_font, args.body_font, args.vietnamese,
        )
    except ComposerError as e:
        print(f"ERROR: {e}", file=sys.stderr)
        return 2
    except Exception as e:  # noqa: BLE001 - top-level CLI boundary, report and exit non-zero
        print(f"ERROR: unexpected failure: {e}", file=sys.stderr)
        return 3

    print(json.dumps(result, indent=2, ensure_ascii=False))
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
