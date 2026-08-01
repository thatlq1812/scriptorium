"""
Shape-role scoring: given a cloned template slide's real shapes, decide
which one is the title, which are body/content candidates.

Ported near-verbatim (pure heuristic scoring, no I/O) from:
D:/elix/archive/platform_archive/modules/presentation/pptx/shape_roles.py
"""
from __future__ import annotations

from typing import Any

import shape_capacity

WATERMARK_TEXT_PATTERNS: tuple[str, ...] = (
    "slidesgo", "freepik", "flaticon", "storyset", "pexels",
    "unsplash", "iconfinder", "icon8", "icon-8",
    "credits:", "credit:", "presentation template",
)

DECORATIVE_TEXT_PATTERNS: tuple[str, ...] = WATERMARK_TEXT_PATTERNS + (
    "click to edit", "click here",
)


def is_watermark_text(text: str | None) -> bool:
    if not text:
        return False
    clean = text.strip().lower()
    if not clean:
        return False
    return any(pattern in clean for pattern in WATERMARK_TEXT_PATTERNS)


def is_decorative_text(text: str | None, top_pct: float, area_pct: float) -> bool:
    if not text:
        return True
    clean = text.strip().lower()
    if not clean:
        return True
    if len(clean) <= 3 and clean.replace(".", "").replace("-", "").isdigit():
        return True
    if is_watermark_text(text):
        return True
    if top_pct > 0.88 and area_pct < 0.05:
        return True
    if area_pct < 0.008:
        return True
    return False


def _group_transform(group_element) -> tuple[int, int, float, float, int, int]:
    """Return (off_x, off_y, scale_x, scale_y, chOff_x, chOff_y) for one
    `<p:grpSp>` element's own `grpSpPr/xfrm` -- the mapping from this
    group's CHILD coordinate space to the space the group itself is
    placed in (its parent's space, or slide space at the top level).
    `shape._element.{x,y,cx,cy,chOff,chExt}` are all get-or-add
    shortcuts (python-pptx never returns None here, defaulting to a
    0/0/0/0 xfrm if the XML genuinely omits one -- verified against a
    real Slidesgo group element before relying on it)."""
    el = group_element
    off_x, off_y = el.x or 0, el.y or 0
    ext_cx, ext_cy = el.cx or 0, el.cy or 0
    ch_off_x, ch_off_y = el.chOff.x or 0, el.chOff.y or 0
    ch_ext_cx, ch_ext_cy = el.chExt.cx or 0, el.chExt.cy or 0
    scale_x = (ext_cx / ch_ext_cx) if ch_ext_cx else 1.0
    scale_y = (ext_cy / ch_ext_cy) if ch_ext_cy else 1.0
    return off_x, off_y, scale_x, scale_y, ch_off_x, ch_off_y


def iter_shapes_with_absolute_bbox(shapes, _transform=None, _depth=0):
    """Recursively walk a shape tree, yielding
    (shape, abs_left, abs_top, abs_width, abs_height, depth) for every
    LEAF shape (recursing into `<p:grpSp>` groups, never yielding the
    group shape itself), with each leaf's bounding box mapped through
    however many nested group transforms sit above it into real,
    absolute SLIDE coordinates.

    Without this, a shape's own `.left/.top/.width/.height` are in its
    immediate parent's coordinate space, not the slide's -- for a
    group whose `chOff/chExt` differs from its own `off/ext` (PowerPoint
    writes this whenever a group is resized after grouping), those raw
    values are meaningfully wrong for area/position-based scoring. Real
    templates surveyed this session have both cases: identity transforms
    (chOff==off, chExt==ext, common in the Slidesgo/Google-Slides
    exports checked) and non-identity ones -- this walker handles both,
    not just the common case, since a heuristic that silently degrades
    on the LESS common real case is exactly the failure mode this
    skill's whole build history (bugs #1-10) kept finding the hard way.
    """
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    for shape in shapes:
        try:
            shape_type = shape.shape_type
        except Exception:
            shape_type = None

        if shape_type == MSO_SHAPE_TYPE.GROUP:
            g_off_x, g_off_y, g_scale_x, g_scale_y, g_ch_off_x, g_ch_off_y = _group_transform(shape._element)
            if _transform is not None:
                p_off_x, p_off_y, p_scale_x, p_scale_y, p_ch_off_x, p_ch_off_y = _transform
                abs_off_x = p_off_x + (g_off_x - p_ch_off_x) * p_scale_x
                abs_off_y = p_off_y + (g_off_y - p_ch_off_y) * p_scale_y
                composed_scale_x = g_scale_x * p_scale_x
                composed_scale_y = g_scale_y * p_scale_y
            else:
                abs_off_x, abs_off_y = g_off_x, g_off_y
                composed_scale_x, composed_scale_y = g_scale_x, g_scale_y
            child_transform = (abs_off_x, abs_off_y, composed_scale_x, composed_scale_y, g_ch_off_x, g_ch_off_y)
            yield from iter_shapes_with_absolute_bbox(shape.shapes, child_transform, _depth + 1)
            continue

        try:
            left, top = int(shape.left or 0), int(shape.top or 0)
            width, height = int(shape.width or 0), int(shape.height or 0)
        except Exception:
            left = top = width = height = 0

        if _transform is not None:
            off_x, off_y, scale_x, scale_y, ch_off_x, ch_off_y = _transform
            abs_left = int(off_x + (left - ch_off_x) * scale_x)
            abs_top = int(off_y + (top - ch_off_y) * scale_y)
            abs_width = int(width * scale_x)
            abs_height = int(height * scale_y)
        else:
            abs_left, abs_top, abs_width, abs_height = left, top, width, height

        yield shape, abs_left, abs_top, abs_width, abs_height, _depth


def analyze_text_shape(
    shape, slide_width_emu: int, slide_height_emu: int,
    bbox: tuple[int, int, int, int] | None = None,
) -> dict[str, Any]:
    """Inspect a text-bearing shape and return scoring metadata.

    `bbox`, when given, is `(left, top, width, height)` in ABSOLUTE
    slide-EMU coordinates already computed by
    `iter_shapes_with_absolute_bbox` (used for a shape nested inside a
    group, where `shape.left/.top/.width/.height` alone would be
    wrong -- see that function's docstring). Omit it for a genuinely
    top-level shape, where `shape.left` etc. already ARE the absolute
    coordinates.
    """
    slide_w = slide_width_emu or 1
    slide_h = slide_height_emu or 1

    text = ""
    para_count = 0
    font_size_pt = 0.0
    try:
        tf = shape.text_frame
        text = (tf.text or "").strip()
        paragraphs = list(tf.paragraphs)
        para_count = len([p for p in paragraphs if (p.text or "").strip()])
        for p in paragraphs:
            for r in p.runs:
                if r.text and r.text.strip() and r.font.size is not None:
                    font_size_pt = float(r.font.size.pt)
                    break
            if font_size_pt > 0:
                break
    except Exception:
        pass

    ph_type = ""
    is_placeholder = False
    try:
        pf = shape.placeholder_format
        if pf is not None and pf.type is not None:
            is_placeholder = True
            ph_type = str(pf.type).split(".")[-1].rstrip("'")
            ph_type = ph_type.split(" ")[0]
    except (ValueError, AttributeError, TypeError):
        pass

    if bbox is not None:
        left, top, width, height = bbox
    else:
        try:
            top = int(shape.top or 0)
            left = int(shape.left or 0)
            width = int(shape.width or 0)
            height = int(shape.height or 0)
        except Exception:
            top = left = width = height = 0

    top_pct = top / slide_h if slide_h else 0.0
    left_pct = left / slide_w if slide_w else 0.0
    width_pct = (width / slide_w) if slide_w else 0.0
    area_pct = (width * height) / (slide_w * slide_h) if slide_w and slide_h else 0.0

    decorative = is_decorative_text(text, top_pct, area_pct)
    watermark = is_watermark_text(text)

    shape_id: int | None = None
    try:
        shape_id = int(shape.shape_id)
    except Exception:
        shape_id = None

    capacity_chars = shape_capacity.estimate_character_capacity(
        width, height, font_size_pt=font_size_pt or 14.0,
    )

    return {
        "shape": shape,
        "shape_id": shape_id,
        "ph_type": ph_type,
        "is_placeholder": is_placeholder,
        "text": text,
        "char_count": len(text),
        "para_count": para_count,
        "top_pct": top_pct,
        "left_pct": left_pct,
        "width_pct": width_pct,
        "area_pct": area_pct,
        "width_emu": width,
        "height_emu": height,
        "font_size_pt": font_size_pt,
        "capacity_chars": capacity_chars,
        "is_decorative": decorative,
        "is_watermark": watermark,
        "is_grouped": False,
    }


def collect_text_shapes(shapes, slide_width_emu: int, slide_height_emu: int) -> list[dict[str, Any]]:
    """Walk the FULL shape tree (recursing into `<p:grpSp>` groups, see
    `iter_shapes_with_absolute_bbox`) and return analyzed text shapes,
    each with a correct absolute-coordinate bounding box regardless of
    nesting depth.

    v0.4.0 change: earlier versions of this skill skipped grouped
    shapes entirely ("mutating text inside a group needs extra
    transform math not needed for the 6 prioritized layouts" -- true
    when written, but real templates surveyed this session (Google-
    Slides-style exports, the dominant real-world source) routinely put
    BOTH decorative art AND real content-bearing text/pictures inside
    groups, so skipping them meant a real, silently-invisible chunk of
    every such template's content -- including candidates for the
    `custom` layout, which depends on the caller being able to see and
    address every real shape on a chosen slide) was never addressable
    at all. The transform math is now implemented (see
    `iter_shapes_with_absolute_bbox`/`_group_transform`), verified
    against a real template's real group nesting before being trusted.
    """
    out: list[dict[str, Any]] = []
    for shape, abs_left, abs_top, abs_width, abs_height, depth in iter_shapes_with_absolute_bbox(shapes):
        try:
            if not shape.has_text_frame:
                continue
        except Exception:
            continue
        info = analyze_text_shape(
            shape, slide_width_emu, slide_height_emu,
            bbox=(abs_left, abs_top, abs_width, abs_height),
        )
        info["group_path"] = []
        info["is_grouped"] = depth > 0
        out.append(info)
    return out


def score_title_shape(info: dict[str, Any]) -> float:
    # v0.4.0: no longer disqualifies a shape purely for being inside a
    # group -- now that iter_shapes_with_absolute_bbox gives grouped
    # shapes correct absolute coordinates, the existing is_decorative/
    # is_watermark checks (already computed against those correct
    # coordinates) are the right filter, not blanket exclusion.
    if info["is_decorative"]:
        return -1
    score = 0.0
    ph = info["ph_type"]
    if ph in ("TITLE", "CENTER_TITLE"):
        score += 1000
    elif ph == "SUBTITLE":
        score += 200
    if info["top_pct"] < 0.45:
        score += 200
    elif info["top_pct"] < 0.60:
        score += 50
    if info["font_size_pt"] >= 30:
        score += 200
    elif info["font_size_pt"] >= 22:
        score += 100
    if 0 < info["char_count"] <= 80:
        score += 100
    elif info["char_count"] <= 150:
        score += 30
    if info["para_count"] == 1:
        score += 50
    if info["width_pct"] > 0.50:
        score += 50
    return score


def score_body_shape(info: dict[str, Any]) -> float:
    """Score a text shape as a candidate to hold this slide's real body
    content.

    Real bug found in a hands-on visual review (rendered to JPG via
    LibreOffice and inspected): the old scoring gave a flat +1000 bonus
    to ANY shape typed BODY/OBJECT/CONTENT, regardless of its actual
    size. A template's 6-item icon-grid slide has 6 tiny CONTENT-typed
    caption placeholders (each maybe 30-40 char capacity) sitting next
    to decorative icon shapes -- one of those tiny captions scored
    highest purely from the placeholder-type bonus and got a whole
    4-bullet paragraph crammed into it, unreadable, while the 6 large
    unrelated icons dominated the visual. Fixed: the placeholder-type
    bonus is now capped well below what a real capacity_chars term can
    swing, so a genuinely spacious BODY placeholder still wins over a
    same-typed tiny one, and a large non-placeholder textbox can
    outscore a tiny placeholder when that's all a template offers.
    """
    if info["is_decorative"]:
        return -1
    score = 0.0
    ph = info["ph_type"]
    if ph in ("BODY", "OBJECT", "CONTENT"):
        score += 150
    elif ph == "SUBTITLE":
        score += 80
    if info["para_count"] >= 3:
        score += 100
    elif info["para_count"] == 2:
        score += 50
    # Dominant term: real character capacity from actual shape
    # dimensions, not just placeholder type or existing text length --
    # capped so one enormous shape can't make every other signal moot.
    score += min(info.get("capacity_chars", 0), 600)
    if 0.20 <= info["top_pct"] <= 0.70:
        score += 60
    if info["area_pct"] >= 0.15:
        score += 100
    elif info["area_pct"] >= 0.06:
        score += 40
    return score


def score_picture_shape(
    shape, slide_width_emu: int, slide_height_emu: int,
    bbox: tuple[int, int, int, int] | None = None,
) -> float:
    """Score a Picture shape as a candidate for the 'image' role of an
    LAYOUT_IMG slide: larger area and more-centered wins. `bbox`
    (absolute left/top/width/height), when given, overrides
    `shape.width`/`.height` -- required for a picture nested inside a
    group, same reasoning as `analyze_text_shape`'s `bbox` param."""
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    try:
        if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
            return -1
    except Exception:
        return -1
    if bbox is not None:
        _left, _top, width, height = bbox
    else:
        try:
            width = int(shape.width or 0)
            height = int(shape.height or 0)
        except Exception:
            return -1
    area_pct = (width * height) / (slide_width_emu * slide_height_emu) if slide_width_emu and slide_height_emu else 0.0
    return area_pct * 1000


def find_all_pictures(shapes) -> list[dict[str, Any]]:
    """Every Picture shape anywhere in the tree (recursing into
    groups, v0.4.0), each with its correct absolute bounding box.
    Returns a list of {"shape": ..., "left"/"top"/"width"/"height":
    absolute EMU, "is_grouped": bool} dicts -- used by
    image_injector.py instead of a top-level-only scan."""
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    out: list[dict[str, Any]] = []
    for shape, left, top, width, height, depth in iter_shapes_with_absolute_bbox(shapes):
        try:
            if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
                continue
        except Exception:
            continue
        out.append({
            "shape": shape, "left": left, "top": top, "width": width, "height": height,
            "is_grouped": depth > 0,
        })
    return out


def find_title_and_body(
    shapes, slide_width_emu: int, slide_height_emu: int
) -> tuple[dict[str, Any] | None, list[dict[str, Any]]]:
    """Return (best_title_info, body_candidate_infos sorted best-first)."""
    infos = collect_text_shapes(shapes, slide_width_emu, slide_height_emu)
    if not infos:
        return None, []

    scored_titles = [(score_title_shape(i), idx) for idx, i in enumerate(infos)]
    scored_titles = [(s, idx) for s, idx in scored_titles if s > 0]
    best_title_idx = None
    if scored_titles:
        scored_titles.sort(key=lambda t: (-t[0], infos[t[1]].get("shape_id") or 0))
        best_title_idx = scored_titles[0][1]

    scored_bodies = [
        (score_body_shape(i), idx) for idx, i in enumerate(infos) if idx != best_title_idx
    ]
    scored_bodies = [(s, idx) for s, idx in scored_bodies if s > 0]
    scored_bodies.sort(key=lambda t: (-t[0], infos[t[1]].get("shape_id") or 0))

    best_title = infos[best_title_idx] if best_title_idx is not None else None
    bodies = [infos[idx] for _s, idx in scored_bodies]
    return best_title, bodies
